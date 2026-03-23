from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿对象
prs = Presentation()

# 定义颜色
COLORS = {
    'primary': RGBColor(102, 126, 234),
    'secondary': RGBColor(118, 75, 162),
    'dark': RGBColor(51, 51, 51),
    'light': RGBColor(245, 247, 250),
    'white': RGBColor(255, 255, 255),
    'gray': RGBColor(153, 153, 153)
}

def add_title_slide(slide, title, subtitle, info=None):
    """添加标题页"""
    # 设置背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary']
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(72)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(32)
    subtitle_para.font.color.rgb = COLORS['white']
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    if info:
        info_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.5))
        info_frame = info_box.text_frame
        info_frame.text = info
        info_para = info_frame.paragraphs[0]
        info_para.font.size = Pt(20)
        info_para.font.color.rgb = COLORS['white']
        info_para.alignment = PP_ALIGN.CENTER

def add_content_slide(slide, title, subtitle=None):
    """添加内容页"""
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.4))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.color.rgb = RGBColor(85, 85, 85)

def add_card(slide, title, content, x, y, width, height):
    """添加卡片"""
    # 添加卡片背景
    shape = slide.shapes.add_shape(1, x, y, width, height)  # 1 = rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['light']
    shape.line.color.rgb = COLORS['primary']
    shape.line.width = Pt(4)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.1), width - Inches(0.4), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['dark']
    
    # 添加内容
    content_box = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.5), width - Inches(0.4), height - Inches(0.6))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    content_frame.text = content
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(16)
    content_para.font.color.rgb = RGBColor(85, 85, 85)

# 幻灯片1：封面页
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_title_slide(slide1, "TRAE", "The Real AI Engineer", "AI原生集成开发环境 · 提升编程效率与质量\n字节跳动推出 · 国内首个AI IDE")

# 幻灯片2：什么是TRAE？
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide2, "什么是TRAE？", "The Real AI Engineer — 真正的AI工程师\nTRAE是由字节跳动推出的国内首个AI原生集成开发环境（AI IDE），旨在通过AI技术提升开发者的编程效率和质量。")

cards2 = [
    ("🎯 核心定位", "能听懂你说话并快速完成代码开发实现的AI助手"),
    ("📅 发布时间", "国内版于2025年3月3日正式发布"),
    ("🤖 模型支持", "搭载豆包1.5 Pro，支持切换满血版DeepSeek R1&V3"),
    ("🌐 官网地址", "国内版：https://www.trae.com.cn/\n国际版：https://www.trae.ai/")
]

for i, (title, content) in enumerate(cards2):
    row = i // 2
    col = i % 2
    add_card(slide2, title, content, Inches(0.5 + col * 4.8), Inches(3.2 + row * 1.8), Inches(4.5), Inches(1.6))

# 幻灯片3：TRAE面向的用户群体
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide3, "TRAE面向的用户群体")

cards3 = [
    ("💻 专业开发者", "提供智能代码补全、Bug修复、代码优化等高级功能，帮助提升编码效率和质量。适用于Web应用开发、工具类应用构建、游戏开发等场景。"),
    ("🎨 非技术背景用户", "通过自然语言交互和Builder模式，无需编程基础即可快速实现项目原型，如定制化游戏、日程管理工具等。"),
    ("🇨🇳 中文开发者", "国内版专为中国开发者优化，提供完整的中文界面、代码注释支持，内置豆包1.5 Pro和DeepSeek R1/V3等本地化模型。"),
    ("🌍 海外开发者", "国际版支持英文界面，集成全球主流模型（如GPT-4o），同时兼容中文输入，满足跨语言开发需求。"),
    ("👥 开发团队与初创企业", "动态协作功能和项目管理工具支持多任务并行处理，帮助团队高效协作，缩短项目周期。")
]

for i, (title, content) in enumerate(cards3):
    if i == 4:  # 最后一个卡片占满一行
        add_card(slide3, title, content, Inches(0.5), Inches(6.8), Inches(9), Inches(1.6))
    else:
        row = i // 2
        col = i % 2
        add_card(slide3, title, content, Inches(0.5 + col * 4.8), Inches(1.5 + row * 1.9), Inches(4.5), Inches(1.7))

# 幻灯片4：核心功能亮点
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide4, "核心功能亮点")

sections4 = [
    ("1️⃣ 动态协作与AI交互", ["Builder模式：通过自然语言描述需求，AI自动生成完整项目框架", "Chat模式：支持实时问答、代码解释、错误修复，提供多模态交互"]),
    ("2️⃣ 智能化编码支持", ["实时代码补全与优化：基于上下文分析，预测并补全代码", "代码片段生成：通过自然语言指令生成跨文件的项目级代码"]),
    ("3️⃣ 多模态与跨平台能力", ["支持图片上传生成代码（如设计草图转前端页面）", "IDE内直接预览Web页面效果", "目前支持macOS，Windows版本正在开发中"])
]

for i, (section_title, items) in enumerate(sections4):
    section_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5 + i * 2.0), Inches(9), Inches(0.5))
    section_frame = section_box.text_frame
    section_frame.text = section_title
    section_para = section_frame.paragraphs[0]
    section_para.font.size = Pt(28)
    section_para.font.bold = True
    section_para.font.color.rgb = COLORS['primary']
    
    for j, item in enumerate(items):
        item_box = slide4.shapes.add_textbox(Inches(0.8), Inches(2.0 + i * 2.0 + j * 0.4), Inches(8.7), Inches(0.35))
        item_frame = item_box.text_frame
        item_frame.text = f"▸ {item}"
        item_para = item_frame.paragraphs[0]
        item_para.font.size = Pt(18)
        item_para.font.color.rgb = RGBColor(85, 85, 85)

# 幻灯片5：核心功能亮点（续）
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide5, "核心功能亮点（续）")

# 添加第4点
section4_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
section4_frame = section4_box.text_frame
section4_frame.text = "4️⃣ 集成主流AI模型"
section4_para = section4_frame.paragraphs[0]
section4_para.font.size = Pt(28)
section4_para.font.bold = True
section4_para.font.color.rgb = COLORS['primary']

content4_box = slide5.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(0.6))
content4_frame = content4_box.text_frame
content4_frame.text = "内置免费模型（豆包1.5 Pro、DeepSeek R1/V3）及国际模型（GPT-4o、Claude-3.5），用户可灵活切换"
content4_para = content4_frame.paragraphs[0]
content4_para.font.size = Pt(18)
content4_para.font.color.rgb = RGBColor(85, 85, 85)

# 添加第5点
section5_box = slide5.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(9), Inches(0.5))
section5_frame = section5_box.text_frame
section5_frame.text = "5️⃣ 开发环境无缝迁移"
section5_para = section5_frame.paragraphs[0]
section5_para.font.size = Pt(28)
section5_para.font.bold = True
section5_para.font.color.rgb = COLORS['primary']

content5_box = slide5.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.6))
content5_frame = content5_box.text_frame
content5_frame.text = "支持从VS Code或Cursor导入配置、插件及快捷键，降低切换成本"
content5_para = content5_frame.paragraphs[0]
content5_para.font.size = Pt(18)
content5_para.font.color.rgb = RGBColor(85, 85, 85)

# 添加特性列表
features5 = [
    ("⚡", "高效协作", "动态协作功能"),
    ("🎯", "智能补全", "上下文感知"),
    ("🖼️", "多模态", "图片转代码"),
    ("🤖", "多模型", "灵活切换"),
    ("🔄", "无缝迁移", "配置导入"),
    ("🌐", "跨平台", "多系统支持")
]

for i, (icon, title, desc) in enumerate(features5):
    col = i % 3
    row = i // 3
    shape = slide5.shapes.add_shape(1, Inches(0.5 + col * 3.2), Inches(4.2 + row * 1.3), Inches(3.0), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['light']
    shape.line.color.rgb = COLORS['primary']
    shape.line.width = Pt(2)
    
    icon_box = slide5.shapes.add_textbox(Inches(1.0 + col * 3.2), Inches(4.3 + row * 1.3), Inches(0.5), Inches(0.4))
    icon_frame = icon_box.text_frame
    icon_frame.text = icon
    icon_para = icon_frame.paragraphs[0]
    icon_para.font.size = Pt(32)
    icon_para.alignment = PP_ALIGN.CENTER
    
    title_box = slide5.shapes.add_textbox(Inches(1.6 + col * 3.2), Inches(4.3 + row * 1.3), Inches(1.9), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(16)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    
    desc_box = slide5.shapes.add_textbox(Inches(0.7 + col * 3.2), Inches(4.7 + row * 1.3), Inches(2.6), Inches(0.5))
    desc_frame = desc_box.text_frame
    desc_frame.text = desc
    desc_para = desc_frame.paragraphs[0]
    desc_para.font.size = Pt(14)
    desc_para.font.color.rgb = RGBColor(85, 85, 85)
    desc_para.alignment = PP_ALIGN.CENTER

# 幻灯片6：重要插件介绍
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide6, "重要插件介绍", "TRAE支持丰富的插件生态，以下是几个核心插件：")

plugins6 = [
    ("🛠️ 代码生成插件", "基于AI的智能代码生成，支持多种编程语言，可根据自然语言描述生成完整代码结构"),
    ("🔍 代码分析插件", "自动分析代码质量，检测潜在问题，提供优化建议，提升代码可维护性"),
    ("📦 依赖管理插件", "智能识别项目依赖，自动添加缺失依赖，更新过时依赖，确保项目环境稳定"),
    ("🌐 版本控制插件", "集成Git等版本控制系统，提供可视化操作界面，简化代码管理流程"),
    ("🧪 测试生成插件", "自动生成单元测试代码，提高测试覆盖率，确保代码质量"),
    ("🎨 UI设计插件", "支持从设计稿生成前端代码，实现设计与开发的无缝衔接")
]

for i, (title, content) in enumerate(plugins6):
    row = i // 2
    col = i % 2
    add_card(slide6, title, content, Inches(0.5 + col * 4.8), Inches(2.1 + row * 1.6), Inches(4.5), Inches(1.4))

# 幻灯片7：插件详细介绍
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide7, "插件详细介绍", "TRAE支持多种实用插件，以下是几个重点推荐：")

plugins7 = [
    ("📄 1. Mintlify 代码转换文档", "将代码自动转换为文档，生成清晰的API文档和使用说明，提高代码可维护性"),
    ("📋 2. Kimi 长文档提炼", "快速总结长文档内容，提取重点信息，生成摘要，提高阅读效率"),
    ("🐛 3. DeepSeek 调试分析", "自动分析代码错误，查找Bug原因，提供修复建议，加速调试过程"),
    ("❌ 4. CodeWhisperer (Amazon)", "注意：此插件目前不可用"),
    ("💡 5. Tabnine", "智能代码补全插件，需购买使用，提供高级代码预测功能")
]

for i, (title, content) in enumerate(plugins7):
    if i == 4:  # 最后一个卡片占满一行
        add_card(slide7, title, content, Inches(0.5), Inches(5.3), Inches(9), Inches(1.4))
    else:
        row = i // 2
        col = i % 2
        add_card(slide7, title, content, Inches(0.5 + col * 4.8), Inches(2.1 + row * 1.6), Inches(4.5), Inches(1.4))

# 幻灯片8：Skills技能
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide8, "Skills技能", "TRAE的Skills功能允许用户创建和使用自定义技能，扩展IDE的能力：")

skills8 = [
    ("🛠️ 创建技能", "通过Builder模式或手动编写代码，创建自定义技能，实现特定功能"),
    ("📦 技能管理", "在技能中心管理已创建的技能，包括启用、禁用、编辑和删除"),
    ("🔄 技能调用", "在对话中显式调用或者对话会自动根据描述调用"),
    ("🌐 技能共享", "将创建的技能分享给团队成员，实现协作开发"),
    ("💡 技能示例pptx", "代码格式化、文档生成、测试用例创建、代码重构等实用技能\n参考：https://github.com/anthropics/skills")
]

for i, (title, content) in enumerate(skills8):
    if i == 4:  # 最后一个卡片占满一行
        add_card(slide8, title, content, Inches(0.5), Inches(5.3), Inches(9), Inches(1.4))
    else:
        row = i // 2
        col = i % 2
        add_card(slide8, title, content, Inches(0.5 + col * 4.8), Inches(2.1 + row * 1.6), Inches(4.5), Inches(1.4))

# 幻灯片9：典型应用场景
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide9, "典型应用场景")

scenarios9 = [
    ("🌐 Web应用开发", "快速生成前后端代码，实现图片上传、数据压缩等功能，并提供界面设计建议"),
    ("🔧 工具类应用", "如图片处理工具、格式转换器等，通过自然语言指令完成核心逻辑开发"),
    ("🎮 游戏开发", "生成贪吃蛇、汉诺塔等小游戏代码，自动处理界面绘制与交互逻辑"),
    ("💻 日常编程辅助", "解释代码含义、优化结构、添加注释，或修复复杂Bug"),
    ("📚 教育与学习", "非技术用户可通过TRAE快速运行开源项目，理解代码逻辑（如翻译界面、添加功能）")
]

for i, (title, content) in enumerate(scenarios9):
    if i == 4:  # 最后一个卡片占满一行
        add_card(slide9, title, content, Inches(0.5), Inches(4.7), Inches(9), Inches(1.4))
    else:
        row = i // 2
        col = i % 2
        add_card(slide9, title, content, Inches(0.5 + col * 4.8), Inches(1.5 + row * 1.6), Inches(4.5), Inches(1.4))

# 幻灯片10：集成的AI模型
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide10, "集成的AI模型")

models10 = [
    ("🇨🇳 国内版模型", ["豆包1.5 Pro", "DeepSeek R1", "DeepSeek V3"], "支持满血版推理能力"),
    ("🌍 国际版模型", ["GPT-4o", "Claude-3.5-Sonnet", "Claude-3.7-Sonnet"], "支持Chat模式、Builder模式、图片理解")
]

for i, (title, tags, desc) in enumerate(models10):
    add_card(slide10, title, desc, Inches(0.5 + i * 4.8), Inches(1.5), Inches(4.5), Inches(2.0))
    
    # 添加标签
    for j, tag in enumerate(tags):
        tag_shape = slide10.shapes.add_shape(1, Inches(0.7 + i * 4.8), Inches(2.0 + j * 0.35), Inches(2.0), Inches(0.3))
        tag_shape.fill.solid()
        tag_shape.fill.fore_color.rgb = COLORS['primary']
        tag_shape.line.color.rgb = COLORS['primary']
        
        tag_box = slide10.shapes.add_textbox(Inches(0.8 + i * 4.8), Inches(2.05 + j * 0.35), Inches(1.8), Inches(0.2))
        tag_frame = tag_box.text_frame
        tag_frame.text = tag
        tag_para = tag_frame.paragraphs[0]
        tag_para.font.size = Pt(14)
        tag_para.font.color.rgb = COLORS['white']
        tag_para.alignment = PP_ALIGN.CENTER

# 添加实操案例
cases_box = slide10.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(0.5))
cases_frame = cases_box.text_frame
cases_frame.text = "📋 实操案例"
cases_para = cases_frame.paragraphs[0]
cases_para.font.size = Pt(28)
cases_para.font.bold = True
cases_para.font.color.rgb = COLORS['primary']

cases10 = [
    ("案例一：井字棋游戏", "使用GLM-4.7，完成\"井字棋\"游戏的代码生成"),
    ("案例二：俄罗斯方块游戏", "通过SOLO模式输入需求，自动生成包含7种经典形状、颜色区分、键盘控制的完整游戏代码")
]

for i, (title, content) in enumerate(cases10):
    case_shape = slide10.shapes.add_shape(1, Inches(0.5), Inches(4.5 + i * 1.3), Inches(9), Inches(1.1))
    case_shape.fill.solid()
    case_shape.fill.fore_color.rgb = RGBColor(249, 249, 249)
    case_shape.line.color.rgb = COLORS['primary']
    case_shape.line.width = Pt(2)
    
    title_box = slide10.shapes.add_textbox(Inches(0.7), Inches(4.6 + i * 1.3), Inches(8.6), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(18)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']
    
    content_box = slide10.shapes.add_textbox(Inches(0.7), Inches(5.0 + i * 1.3), Inches(8.6), Inches(0.5))
    content_frame = content_box.text_frame
    content_frame.text = content
    content_para = content_frame.paragraphs[0]
    content_para.font.size = Pt(16)
    content_para.font.color.rgb = RGBColor(85, 85, 85)

# 幻灯片11：总结
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_content_slide(slide11, "总结", "TRAE作为国内首个AI原生集成开发环境，通过以下核心优势为开发者赋能：")

summary11 = [
    "智能化：AI驱动的代码生成、补全、优化和调试",
    "易用性：自然语言交互，降低编程门槛",
    "多模态：支持文字、图片等多种输入方式",
    "灵活性：多模型支持，满足不同需求",
    "协作性：团队协作功能，提升开发效率"
]

for i, item in enumerate(summary11):
    item_box = slide11.shapes.add_textbox(Inches(0.8), Inches(2.2 + i * 0.5), Inches(8.7), Inches(0.4))
    item_frame = item_box.text_frame
    item_frame.text = f"▸ {item}"
    item_para = item_frame.paragraphs[0]
    item_para.font.size = Pt(18)
    item_para.font.color.rgb = RGBColor(85, 85, 85)

# 添加结束卡片
end_card = slide11.shapes.add_shape(1, Inches(0.5), Inches(5.0), Inches(9), Inches(1.5))
end_card.fill.solid()
end_card.fill.fore_color.rgb = COLORS['light']
end_card.line.color.rgb = COLORS['primary']
end_card.line.width = Pt(4)

end_title_box = slide11.shapes.add_textbox(Inches(0.7), Inches(5.2), Inches(8.6), Inches(0.5))
end_title_frame = end_title_box.text_frame
end_title_frame.text = "🚀 开始使用TRAE，开启AI编程新时代！"
end_title_para = end_title_frame.paragraphs[0]
end_title_para.font.size = Pt(28)
end_title_para.font.bold = True
end_title_para.font.color.rgb = COLORS['dark']
end_title_para.alignment = PP_ALIGN.CENTER

end_info_box = slide11.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(8.6), Inches(0.6))
end_info_frame = end_info_box.text_frame
end_info_frame.text = "国内版：https://www.trae.com.cn/\n国际版：https://www.trae.ai/"
end_info_para = end_info_frame.paragraphs[0]
end_info_para.font.size = Pt(18)
end_info_para.font.color.rgb = RGBColor(85, 85, 85)
end_info_para.alignment = PP_ALIGN.CENTER

# 保存PPT文件
prs.save('trae_training_presentation.pptx')
print("PPT文件已成功生成: trae_training_presentation.pptx")